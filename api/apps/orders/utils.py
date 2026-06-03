import re
import unicodedata
import zipfile
import docx
import pypdf
import os
import subprocess
import tempfile
from io import BytesIO

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import xlrd
    HAS_XLRD = True
except ImportError:
    HAS_XLRD = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def _safe_decode(data: bytes) -> str:
    if data.startswith(b'\xff\xfe') or data.startswith(b'\xfe\xff'):
        return data.decode('utf-16', errors='replace')
    if data.startswith(b'\xef\xbb\xbf'):
        return data.decode('utf-8-sig', errors='replace')
    return data.decode('utf-8', errors='replace')


def normalize_text(text: str, collapse_spaces: bool = True) -> str:
    text = text.replace('\x02', '')
    text = text.replace('\x07', '')
    text = text.replace('\x08', '')
    text = text.replace('\x0b', '\n')
    text = text.replace('\x0c', '\n')
    text = text.replace('\x13', '')
    text = text.replace('\x14', '')
    text = text.replace('\x15', '')
    text = text.replace('\r', '')
    text = text.replace('\xad', '')
    text = text.replace('\ufeff', '')
    text = text.replace('\u200b', '')
    text = text.replace('\u200c', '')
    text = text.replace('\u200d', '')
    text = text.replace('\u2028', '\n')
    text = text.replace('\u2029', '\n')

    spaces_to_normalize = [
        '\t', '\xa0', '\u2000', '\u2001', '\u2002', '\u2003', '\u2004',
        '\u2005', '\u2006', '\u2007', '\u2008', '\u2009', '\u200a',
        '\u202f', '\u205f', '\u3000'
    ]
    for space in spaces_to_normalize:
        text = text.replace(space, ' ')

    text = re.sub(r'\n{2,}', '\n', text)
    if collapse_spaces:
        text = re.sub(r' {2,}', ' ', text)

    text = unicodedata.normalize('NFC', text)
    return text.strip()


def extract_text_from_shapes(doc) -> list[str]:
    extra_texts = []
    try:
        for txbx in doc.element.xpath('.//w:txbxContent'):
            for p in txbx.xpath('.//w:p'):
                texts = [node.text for node in p.xpath('.//w:t') if node.text]
                if texts:
                    extra_texts.append("".join(texts))
    except Exception:
        pass

    try:
        for section in doc.sections:
            for header in [section.header, section.first_page_header, section.even_page_header]:
                if header and not header.is_linked_to_previous:
                    for p in header.paragraphs:
                        if p.text: extra_texts.append(p.text)
                    for table in header.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text: extra_texts.append(cell.text)

            for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
                if footer and not footer.is_linked_to_previous:
                    for p in footer.paragraphs:
                        if p.text: extra_texts.append(p.text)
                    for table in footer.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text: extra_texts.append(cell.text)
    except Exception:
        pass

    return extra_texts


def _get_docx_pages_from_metadata(file) -> int:
    """Читає точну кількість сторінок з docProps/app.xml."""
    try:
        file.seek(0)
        if zipfile.is_zipfile(file):
            file.seek(0)
            with zipfile.ZipFile(file) as archive:
                if 'docProps/app.xml' in archive.namelist():
                    app_xml = archive.read('docProps/app.xml').decode('utf-8', errors='replace')
                    match = re.search(r'<Pages>(\d+)</Pages>', app_xml)
                    if match:
                        return int(match.group(1))
    except Exception:
        pass
    return 0


def _antiword_extract(file) -> str:
    try:
        file.seek(0)
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp:
            tmp.write(file.read())
            tmp_path = tmp.name
        result = subprocess.run(['antiword', tmp_path], capture_output=True, text=True, timeout=15)
        os.unlink(tmp_path)
        if result.returncode == 0:
            return result.stdout
    except Exception:
        pass
    return ''


def _extract_xliff_text(data: bytes) -> list[str]:
    parts = []
    try:
        text = data.decode('utf-8', errors='replace')
        for tag in ('target', 'seg-source', 'source'):
            for match in re.finditer(rf'<{tag}[^>]*>(.*?)</{tag}>', text, re.DOTALL):
                inner = match.group(1)
                clean = re.sub(r'<[^>]+>', '', inner).strip()
                if clean:
                    parts.append(clean)
            if parts:
                break
    except Exception:
        pass
    return parts


def _extract_tmx_text(data: bytes) -> list[str]:
    parts = []
    try:
        text = _safe_decode(data)
        for match in re.finditer(r'<seg[^>]*>(.*?)</seg>', text, re.IGNORECASE | re.DOTALL):
            clean = re.sub(r'<[^>]+>', '', match.group(1)).strip()
            if clean:
                parts.append(clean)
    except Exception:
        pass
    return parts


def _extract_ttx_text(data: bytes) -> list[str]:
    parts = []
    try:
        text = data.decode('utf-8', errors='replace')
        for match in re.finditer(r'<Tu\b[^>]*>(.*?)</Tu>', text, re.DOTALL):
            clean = re.sub(r'<[^>]+>', '', match.group(1)).strip()
            if clean:
                parts.append(clean)
    except Exception:
        pass
    return parts


def _extract_mqxlz(file) -> tuple[list[str], int]:
    parts = []
    doc_count = 0
    try:
        file.seek(0)
        if zipfile.is_zipfile(file):
            file.seek(0)
            with zipfile.ZipFile(file) as zf:
                xliff_files = [n for n in zf.namelist() if n.endswith('.mqxliff')]
                doc_count = len(xliff_files)
                for name in xliff_files:
                    parts.extend(_extract_xliff_text(zf.read(name)))
    except Exception:
        pass
    return parts, doc_count


def _extract_sdltm(file) -> tuple[list[str], int]:
    parts = []
    pages = 0
    tmp_path = None
    try:
        file.seek(0)
        raw_data = file.read()
        if not raw_data:
            return [], 0
        with tempfile.NamedTemporaryFile(suffix='.sdltm', delete=False) as tmp:
            tmp.write(raw_data)
            tmp_path = tmp.name
        import sqlite3
        conn = sqlite3.connect(tmp_path)
        cursor = conn.cursor()
        try:
            cursor.execute("SELECT COUNT(*) FROM translation_units")
            pages = max(1, cursor.fetchone()[0] // 250)
        except Exception:
            pages = 1
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = [t[0] for t in cursor.fetchall()]
        for table in tables:
            try:
                cursor.execute(f"SELECT * FROM {table}")
                for row in cursor.fetchall():
                    for item in row:
                        val = ""
                        if isinstance(item, str):
                            val = item
                        elif isinstance(item, bytes):
                            val = item.decode('utf-8', errors='ignore')
                        if val and len(val) > 2:
                            clean = re.sub(r'<[^>]+>', ' ', val).strip()
                            if clean and len(clean) > 2 and not clean.isnumeric():
                                parts.append(clean)
            except Exception:
                continue
        conn.close()
    except Exception as e:
        print(f"Critical SDLTm Error: {e}")
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except Exception:
                pass
    return parts, pages


def analyze_file_content(file):
    stats = {
        "chars_with_spaces": 0,
        "chars_no_spaces": 0,
        "pages": 0,
        "images": 0
    }
    file_name = file.name.lower()
    text_parts = []
    stats_from_metadata = False
    file.seek(0)

    try:
        # ------------------------------------------------------------------ #
        # WORD (.docx / .doc / .dotx / .dotm / .docm)
        # ------------------------------------------------------------------ #
        if file_name.endswith(('.docx', '.doc', '.dotx', '.dotm', '.docm')):
            file.seek(0)
            try:
                if zipfile.is_zipfile(file):
                    file.seek(0)
                    with zipfile.ZipFile(file) as archive:
                        namelist = archive.namelist()
                        stats["images"] = sum(1 for name in namelist if name.startswith('word/media/'))
            except Exception:
                pass

            # Читаємо точну кількість сторінок з метаданих (тільки для zip-форматів)
            if file_name.endswith(('.docx', '.dotx', '.dotm', '.docm')):
                pages_from_meta = _get_docx_pages_from_metadata(file)
                if pages_from_meta > 0:
                    stats["pages"] = pages_from_meta

            try:
                file.seek(0)
                doc = docx.Document(file)

                for p in doc.paragraphs:
                    if p.text:
                        is_list = False
                        if p.style and p.style.name and 'List' in p.style.name:
                            is_list = True
                        elif p._element.xpath('.//w:numPr'):
                            is_list = True

                        if is_list:
                            text_parts.append("· " + p.text)
                        else:
                            text_parts.append(p.text)

                visited_cells = set()
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell not in visited_cells:
                                visited_cells.add(cell)
                                if cell.text:
                                    text_parts.append(cell.text)

                text_parts.extend(extract_text_from_shapes(doc))

            except Exception:
                pass

            if not text_parts and file_name.endswith('.doc'):
                extracted = _antiword_extract(file)
                if extracted.strip():
                    text_parts.append(extracted)

            # Fallback для .doc або якщо метадані не дали результату
            if stats["pages"] == 0 and text_parts:
                full_text_preview = normalize_text("\n".join(p for p in text_parts if p))
                char_count = len(full_text_preview.replace('\n', '').replace(' ', ''))
                stats["pages"] = max(1, round(char_count / 1800))

        # ------------------------------------------------------------------ #
        # PDF
        # ------------------------------------------------------------------ #
        elif file_name.endswith('.pdf'):
            try:
                file.seek(0)
                reader = pypdf.PdfReader(file)
                stats["pages"] = len(reader.pages)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        cleaned = re.sub(r'\n+', ' ', extracted)
                        cleaned = re.sub(r' {2,}', ' ', cleaned)
                        text_parts.append(cleaned.strip())
                    if hasattr(page, 'images'):
                        stats["images"] += len(page.images)
            except Exception:
                pass

        # ------------------------------------------------------------------ #
        # SDLTM
        # ------------------------------------------------------------------ #
        elif file_name.endswith('.sdltm'):
            try:
                parts, pages = _extract_sdltm(file)
                if parts:
                    text_parts.extend(parts)
                if pages > 0:
                    stats["pages"] = max(stats.get("pages", 0), pages)
            except Exception:
                pass

        # ------------------------------------------------------------------ #
        # EXCEL
        # ------------------------------------------------------------------ #
        elif file_name.endswith(('.xlsx', '.xlsm', '.xltx', '.xltm')):
            if HAS_OPENPYXL:
                try:
                    file.seek(0)
                    wb = openpyxl.load_workbook(file, read_only=True, data_only=True)
                    stats["pages"] = len(wb.sheetnames)
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            for cell in row:
                                if cell is not None:
                                    text_parts.append(str(cell))
                    wb.close()
                except Exception:
                    pass
            file.seek(0)
            try:
                if zipfile.is_zipfile(file):
                    file.seek(0)
                    with zipfile.ZipFile(file) as archive:
                        stats["images"] = sum(1 for name in archive.namelist() if name.startswith('xl/media/'))
            except Exception:
                pass

        elif file_name.endswith('.xls'):
            if HAS_XLRD:
                try:
                    file.seek(0)
                    raw = file.read()
                    wb = xlrd.open_workbook(file_contents=raw)
                    stats["pages"] = wb.nsheets
                    for sheet in wb.sheets():
                        for row_idx in range(sheet.nrows):
                            for col_idx in range(sheet.ncols):
                                val = sheet.cell(row_idx, col_idx).value
                                if val not in (None, ''):
                                    text_parts.append(str(val))
                except Exception:
                    pass

        # ------------------------------------------------------------------ #
        # POWERPOINT
        # ------------------------------------------------------------------ #
        elif file_name.endswith(('.pptx', '.pptm', '.ppsx', '.potx')):
            if HAS_PPTX:
                try:
                    file.seek(0)
                    prs = Presentation(file)
                    stats["pages"] = len(prs.slides)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    text = para.text.strip()
                                    if text:
                                        text_parts.append(text)
                            if shape.has_table:
                                for row in shape.table.rows:
                                    for cell in row.cells:
                                        if cell.text.strip():
                                            text_parts.append(cell.text.strip())
                except Exception:
                    pass
            file.seek(0)
            try:
                if zipfile.is_zipfile(file):
                    file.seek(0)
                    with zipfile.ZipFile(file) as archive:
                        stats["images"] = sum(1 for name in archive.namelist() if name.startswith('ppt/media/'))
            except Exception:
                pass

        # ------------------------------------------------------------------ #
        # ПЕРЕКЛАДАЦЬКІ ФОРМАТИ (Trados, memoQ, etc.)
        # ------------------------------------------------------------------ #
        elif file_name.endswith('.sdlxliff'):
            try:
                file.seek(0)
                text_parts.extend(_extract_xliff_text(file.read()))
                stats["pages"] = 1
            except Exception:
                pass

        elif file_name.endswith(('.sdlppx', '.sdlproj', '.wsxz')):
            try:
                file.seek(0)
                if zipfile.is_zipfile(file):
                    file.seek(0)
                    with zipfile.ZipFile(file) as zf:
                        xliff_files = [n for n in zf.namelist() if n.endswith(('.sdlxliff', '.xliff', '.xlf'))]
                        xlz_files = [n for n in zf.namelist() if n.endswith('.xlz')]
                        stats["pages"] = (len(xliff_files) + len(xlz_files)) or 1
                        for name in xliff_files:
                            text_parts.extend(_extract_xliff_text(zf.read(name)))
                        for name in xlz_files:
                            xlz_data = zf.read(name)
                            with zipfile.ZipFile(BytesIO(xlz_data)) as inner_zf:
                                inner_xlf = [n for n in inner_zf.namelist() if n.endswith(('.xlf', '.xliff'))]
                                for inner_name in inner_xlf:
                                    text_parts.extend(_extract_xliff_text(inner_zf.read(inner_name)))
            except Exception:
                pass

        elif file_name.endswith('.ttx'):
            try:
                file.seek(0)
                text_parts.extend(_extract_ttx_text(file.read()))
                stats["pages"] = 1
            except Exception:
                pass

        elif file_name.endswith('.tmx'):
            try:
                file.seek(0)
                raw_data = file.read()
                text_parts.extend(_extract_tmx_text(raw_data))
                content = _safe_decode(raw_data)
                tu_count = len(re.findall(r'<tu\b', content, re.IGNORECASE))
                stats["pages"] = max(1, tu_count // 250)
            except Exception:
                pass

        elif file_name.endswith('.mqxliff'):
            try:
                file.seek(0)
                text_parts.extend(_extract_xliff_text(file.read()))
                stats["pages"] = 1
            except Exception:
                pass

        elif file_name.endswith('.mqxlz'):
            parts, doc_count = _extract_mqxlz(file)
            text_parts.extend(parts)
            stats["pages"] = max(1, doc_count)

        elif file_name.endswith(('.xliff', '.xlf')):
            try:
                file.seek(0)
                text_parts.extend(_extract_xliff_text(file.read()))
                stats["pages"] = 1
            except Exception:
                pass

        # ------------------------------------------------------------------ #
        # РАХУЄМО СИМВОЛИ
        # ------------------------------------------------------------------ #
        if not stats_from_metadata:
            full_text = "\n".join(p for p in text_parts if p)

            needs_space_collapse = not file_name.endswith(
                ('.docx', '.doc', '.xlsx', '.xls', '.pptx', '.xliff', '.mqxliff', '.xlf', '.sdlxliff'))

            full_text = normalize_text(full_text, collapse_spaces=needs_space_collapse)
            text_for_counting = full_text.replace('\n', '')

            stats["chars_with_spaces"] = len(text_for_counting)
            stats["chars_no_spaces"] = len(text_for_counting.replace(' ', ''))

    except Exception as e:
        print(f"Error analyzing file {file_name}: {e}")

    file.seek(0)
    return stats